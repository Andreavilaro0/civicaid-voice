#!/usr/bin/env python3
"""Generate Clara presentation PPTX."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──
TRUST   = RGBColor(0x1B, 0x5E, 0x7B)
HOPE    = RGBColor(0x2E, 0x7D, 0x4F)
WARMTH  = RGBColor(0xD4, 0x6A, 0x1E)
ALERT   = RGBColor(0xC6, 0x28, 0x28)
DARK    = RGBColor(0x0A, 0x0A, 0x0A)
WHITE   = RGBColor(0xFA, 0xFA, 0xF9)
MUTED   = RGBColor(0x75, 0x75, 0x75)
LIGHT50 = RGBColor(0x99, 0x99, 0x99)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height,
             font_size=18, color=WHITE, bold=False, italic=False,
             alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_shape_text(slide, text, left, top, width, height,
                   fill_color=None, border_color=None,
                   font_size=14, font_color=WHITE, bold=False,
                   alignment=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = alignment
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = "Calibri"
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


# ════════════════════════════════════════════
# SLIDE 1: COVER
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

add_text(slide, "Andrea Avila", Inches(0), Inches(0.5), W, Inches(0.5),
         font_size=14, color=LIGHT50, alignment=PP_ALIGN.CENTER)

# Logo circles (simplified)
cx, cy = Inches(6.2), Inches(1.8)
add_circle(slide, cx, cy, Inches(0.5), WARMTH)

add_text(slide, "Clara", Inches(0), Inches(2.8), W, Inches(1.2),
         font_size=72, color=WHITE, bold=False, alignment=PP_ALIGN.CENTER,
         font_name="Georgia")

add_text(slide, "Asistente conversacional que ayuda a personas\nvulnerables a navegar ayudas del gobierno",
         Inches(2), Inches(4.1), Inches(9), Inches(1),
         font_size=24, color=LIGHT50, alignment=PP_ALIGN.CENTER)

# Pills
pills = ["WhatsApp-first", "Voz + Texto + Foto", "ES / FR"]
for i, pill in enumerate(pills):
    x = Inches(4 + i * 2)
    add_shape_text(slide, pill, x, Inches(5.3), Inches(1.8), Inches(0.45),
                   border_color=LIGHT50, font_size=11, font_color=LIGHT50)

add_text(slide, "OdiseIA4Good  ·  UDIT  ·  Febrero 2026",
         Inches(0), Inches(6.5), W, Inches(0.4),
         font_size=12, color=LIGHT50, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════
# SLIDE 2: EL PROBLEMA
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, "EL PROBLEMA", Inches(0), Inches(0.5), W, Inches(0.4),
         font_size=12, color=LIGHT50, alignment=PP_ALIGN.CENTER)

add_text(slide, "11.4M", Inches(0), Inches(1.3), W, Inches(1),
         font_size=96, color=WARMTH, bold=False, alignment=PP_ALIGN.CENTER,
         font_name="Georgia")

add_text(slide, "personas en riesgo de exclusion social en Espana",
         Inches(2), Inches(2.5), Inches(9), Inches(0.6),
         font_size=22, color=LIGHT50, alignment=PP_ALIGN.CENTER)

# Stats row
stats = [
    ("73%", "no completa tramites\npor barreras digitales", ALERT),
    ("40%", "no habla espanol\ncomo lengua materna", TRUST),
    ("87%", "usa WhatsApp\na diario", HOPE),
]
for i, (num, label, color) in enumerate(stats):
    x = Inches(2 + i * 3.2)
    add_text(slide, num, x, Inches(3.5), Inches(2.5), Inches(0.8),
             font_size=52, color=color, alignment=PP_ALIGN.CENTER,
             font_name="Georgia")
    add_text(slide, label, x, Inches(4.5), Inches(2.5), Inches(0.8),
             font_size=14, color=LIGHT50, alignment=PP_ALIGN.CENTER)

add_text(slide, "Las ayudas existen. Entre el papel y la persona hay un muro.",
         Inches(2), Inches(6), Inches(9), Inches(0.5),
         font_size=18, color=LIGHT50, italic=True, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════
# SLIDE 3: PERSONAS
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, "PARA QUIEN ES CLARA", Inches(0), Inches(0.4), W, Inches(0.4),
         font_size=12, color=LIGHT50, alignment=PP_ALIGN.CENTER)

personas = [
    ("Maria, 74 anos", "Le tiemblan las manos.\nNo puede teclear.\nQuiere saber del IMV.", "NOTA DE VOZ", WARMTH),
    ("Ahmed, 28 anos", "Habla frances.\nNecesita empadronarse.\nTodo esta en espanol.", "BILINGUE ES/FR", TRUST),
    ("Fatima, 45 anos", "Habla pero no lee.\nNecesita tarjeta sanitaria\npara su hijo.", "AUDIO RESPUESTA", HOPE),
]

for i, (name, desc, tag, color) in enumerate(personas):
    x = Inches(1 + i * 3.8)
    # Card bg
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, Inches(1.2), Inches(3.3), Inches(3.8))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x18, 0x18, 0x18)
    card.line.color.rgb = RGBColor(0x30, 0x30, 0x30)
    card.line.width = Pt(1)

    # Icon circle
    add_circle(slide, Inches(1.3 + i * 3.8), Inches(1.5), Inches(0.6), color)

    add_text(slide, name, Inches(1.3 + i * 3.8), Inches(2.3), Inches(2.8), Inches(0.4),
             font_size=18, color=WHITE, bold=True)
    add_text(slide, desc, Inches(1.3 + i * 3.8), Inches(2.8), Inches(2.8), Inches(1.2),
             font_size=14, color=LIGHT50)

    # Tag
    add_shape_text(slide, tag, Inches(1.3 + i * 3.8), Inches(4.2), Inches(1.6), Inches(0.35),
                   fill_color=RGBColor(0x20, 0x20, 0x20), border_color=color,
                   font_size=9, font_color=color, bold=True)

add_text(slide, "Clara no pide que las personas se adapten. Clara se adapta a las personas.",
         Inches(1), Inches(5.5), Inches(11), Inches(0.5),
         font_size=18, color=WARMTH, italic=True, alignment=PP_ALIGN.CENTER)

# More vulnerable people
more = "Tambien: mayores solos · inmigrantes sin idioma · mujeres en violencia de genero · personas sin hogar · discapacidad cognitiva · analfabetismo funcional"
add_text(slide, more, Inches(1.5), Inches(6.2), Inches(10), Inches(0.6),
         font_size=12, color=LIGHT50, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════
# SLIDE 4: USER JOURNEY
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, "EJEMPLO: MARIA QUIERE SOLICITAR EL IMV",
         Inches(0), Inches(0.4), W, Inches(0.4),
         font_size=12, color=LIGHT50, alignment=PP_ALIGN.CENTER)

steps = [
    ("1", "Abre WhatsApp", "Lo que ya tiene\nen el movil", WARMTH),
    ("2", "Nota de voz", '"Quiero saber\ndel IMV"', HOPE),
    ("3", "Whisper + Gemini", "Transcribe,\nentiende, busca", TRUST),
    ("4", "Respuesta + audio", "Texto claro +\nnota de voz", HOPE),
    ("5", "Sabe que hacer", "Paso concreto\npara manana", HOPE),
]

# Timeline line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(1.5), Inches(3.1), Inches(10.3), Pt(2))
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(0x30, 0x30, 0x30)
line.line.fill.background()

for i, (num, label, desc, color) in enumerate(steps):
    x = Inches(1.5 + i * 2.2)

    # Circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  x + Inches(0.35), Inches(2.2), Inches(0.7), Inches(0.7))
    circ.fill.solid()
    circ.fill.fore_color.rgb = RGBColor(0x15, 0x15, 0x15)
    circ.line.color.rgb = color
    circ.line.width = Pt(2)

    tf = circ.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = color
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text(slide, label, x, Inches(3.3), Inches(1.5), Inches(0.4),
             font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, desc, x, Inches(3.8), Inches(1.5), Inches(0.8),
             font_size=12, color=LIGHT50, alignment=PP_ALIGN.CENTER)

add_text(slide, "Todo en < 30 segundos. Sin apps. Sin formularios. Sin esperas.",
         Inches(2), Inches(5.5), Inches(9), Inches(0.5),
         font_size=16, color=LIGHT50, alignment=PP_ALIGN.CENTER)

# More use cases
add_text(slide, "Otros usos: empadronamiento · tarjeta sanitaria · NIE/TIE · reagrupacion familiar · ayudas vivienda · becas",
         Inches(1.5), Inches(6.2), Inches(10), Inches(0.5),
         font_size=12, color=LIGHT50, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════
# SLIDE 5: DEMO WHATSAPP
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, "DEMO", Inches(0.8), Inches(0.5), Inches(4), Inches(0.4),
         font_size=12, color=LIGHT50)

add_text(slide, "Clara en\nWhatsApp", Inches(0.8), Inches(1.2), Inches(5), Inches(1.5),
         font_size=48, color=WHITE, font_name="Georgia")

# Features
features = [
    "Texto, voz o foto de documento",
    "Espanol y frances automatico",
    "Audio para quienes no leen",
    "Guardrails de seguridad",
    "Fuentes oficiales verificadas",
]
for i, feat in enumerate(features):
    add_text(slide, "✓  " + feat, Inches(0.8), Inches(3 + i * 0.45), Inches(5), Inches(0.4),
             font_size=15, color=LIGHT50)

# Phone mockup placeholder
phone = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(7.5), Inches(0.5), Inches(4.5), Inches(6.5))
phone.fill.solid()
phone.fill.fore_color.rgb = RGBColor(0x1C, 0x1C, 0x1E)
phone.line.color.rgb = RGBColor(0x3A, 0x3A, 0x3C)
phone.line.width = Pt(2)

# Phone header
header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(7.7), Inches(0.7), Inches(4.1), Inches(0.6))
header.fill.solid()
header.fill.fore_color.rgb = TRUST
header.line.fill.background()
add_text(slide, "  Clara                                    ES | FR",
         Inches(7.7), Inches(0.7), Inches(4.1), Inches(0.6),
         font_size=13, color=WHITE, bold=True)

# Chat bubbles
bubbles = [
    (False, "Hola, soy Clara. Estoy aqui para\nayudarte con tramites y ayudas."),
    (True, "Quiero saber que es el IMV"),
    (False, "El IMV es una ayuda economica\npara personas con pocos recursos.\nPuede ser 604-1.200€/mes.\n\n▶ Escuchar respuesta  0:34"),
    (True, "🎙 Nota de voz 0:05"),
    (False, "Para solicitarlo necesitas:\n1. DNI/NIE\n2. Empadronamiento\n3. Declaracion renta"),
]

y = 1.5
for is_user, text in bubbles:
    lines = text.count('\n') + 1
    h = 0.2 + lines * 0.22
    if is_user:
        x = Inches(9.5)
        w = Inches(2.1)
        bg_color = TRUST
        fc = WHITE
    else:
        x = Inches(7.9)
        w = Inches(3)
        bg_color = RGBColor(0xE3, 0xF2, 0xFD)
        fc = RGBColor(0x1A, 0x1A, 0x1A)

    bub = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 x, Inches(y), w, Inches(h))
    bub.fill.solid()
    bub.fill.fore_color.rgb = bg_color
    bub.line.fill.background()

    tf = bub.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.color.rgb = fc
    tf.paragraphs[0].font.name = "Calibri"

    y += h + 0.12

add_text(slide, "→ Abrir demo-whatsapp.html para demo interactiva",
         Inches(7.5), Inches(6.3), Inches(4.5), Inches(0.4),
         font_size=10, color=LIGHT50, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════
# SLIDE 6: DEMO WEB APP
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, "INTERFAZ ACCESIBLE", Inches(7), Inches(0.5), Inches(5), Inches(0.4),
         font_size=12, color=LIGHT50)

add_text(slide, "Clara Web", Inches(7), Inches(1.2), Inches(5), Inches(0.8),
         font_size=48, color=WHITE, font_name="Georgia")

add_text(slide, "Para trabajadoras sociales, ONGs\ny mediadores comunitarios.",
         Inches(7), Inches(2.2), Inches(5), Inches(0.7),
         font_size=18, color=LIGHT50)

web_features = [
    "WCAG AA accesible",
    "Grabacion de voz integrada",
    "Mobile-first",
    "Selector de idioma ES / FR",
]
for i, feat in enumerate(web_features):
    add_text(slide, "✓  " + feat, Inches(7), Inches(3.2 + i * 0.45), Inches(5), Inches(0.4),
             font_size=15, color=LIGHT50)

# Phone mockup (left side)
phone2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(1), Inches(0.5), Inches(4.5), Inches(6.5))
phone2.fill.solid()
phone2.fill.fore_color.rgb = RGBColor(0x1C, 0x1C, 0x1E)
phone2.line.color.rgb = RGBColor(0x3A, 0x3A, 0x3C)
phone2.line.width = Pt(2)

header2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(1.2), Inches(0.7), Inches(4.1), Inches(0.6))
header2.fill.solid()
header2.fill.fore_color.rgb = TRUST
header2.line.fill.background()
add_text(slide, "  Clara Web                              ES",
         Inches(1.2), Inches(0.7), Inches(4.1), Inches(0.6),
         font_size=13, color=WHITE, bold=True)

web_bubbles = [
    (False, "Bienvenido/a. Soy Clara,\ntu asistente para tramites."),
    (True, "Como pido la tarjeta sanitaria?"),
    (False, "Para la tarjeta sanitaria:\n\n✅ Empadronamiento\n✅ Pasaporte o NIE\n✅ Centro de salud\n\n▶ Escuchar  0:28"),
]

y = 1.5
for is_user, text in web_bubbles:
    lines = text.count('\n') + 1
    h = 0.2 + lines * 0.22
    if is_user:
        x = Inches(3.2)
        w = Inches(2.1)
        bg_color = TRUST
        fc = WHITE
    else:
        x = Inches(1.4)
        w = Inches(3.2)
        bg_color = RGBColor(0xE3, 0xF2, 0xFD)
        fc = RGBColor(0x1A, 0x1A, 0x1A)

    bub = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 x, Inches(y), w, Inches(h))
    bub.fill.solid()
    bub.fill.fore_color.rgb = bg_color
    bub.line.fill.background()

    tf = bub.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.color.rgb = fc

    y += h + 0.12

add_text(slide, "→ Abrir demo-webapp.html para demo interactiva",
         Inches(1), Inches(6.3), Inches(4.5), Inches(0.4),
         font_size=10, color=LIGHT50, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════
# SLIDE 7: ARQUITECTURA
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, "COMO FUNCIONA", Inches(0), Inches(0.3), W, Inches(0.4),
         font_size=12, color=LIGHT50, alignment=PP_ALIGN.CENTER)
add_text(slide, "Arquitectura", Inches(0), Inches(0.8), W, Inches(0.7),
         font_size=36, color=WHITE, alignment=PP_ALIGN.CENTER, font_name="Georgia")

# Architecture nodes
arch_nodes = [
    ("WhatsApp\n+ Web", "Texto / Voz / Foto", HOPE),
    ("Twilio", "Mensajeria\nsegura", TRUST),
    ("Whisper", "Voz → Texto\nOpenAI", WARMTH),
    ("Gemini\n1.5 Flash", "IA generativa\nRAG + KB", TRUST),
    ("Respuesta\nverificada", "Texto + Audio\n+ Fuentes", HOPE),
]

for i, (name, desc, color) in enumerate(arch_nodes):
    x = Inches(0.8 + i * 2.5)
    # Box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 x, Inches(2.2), Inches(1.8), Inches(1.4))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x15, 0x15, 0x15)
    box.line.color.rgb = color
    box.line.width = Pt(1.5)

    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text(slide, desc, x, Inches(3.8), Inches(1.8), Inches(0.6),
             font_size=10, color=LIGHT50, alignment=PP_ALIGN.CENTER)

    # Arrow
    if i < len(arch_nodes) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       x + Inches(1.9), Inches(2.75), Inches(0.5), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0x40, 0x40, 0x40)
        arrow.line.fill.background()

# Tech stack pills
techs = ["Python + Flask", "Docker", "Render", "PostgreSQL + pgvector",
         "React + TypeScript", "469 tests", "11 skills", "8 KBs"]
for i, tech in enumerate(techs):
    col = i % 4
    row = i // 4
    x = Inches(2 + col * 2.4)
    y = Inches(5.0 + row * 0.5)
    add_shape_text(slide, tech, x, y, Inches(2), Inches(0.38),
                   border_color=RGBColor(0x40, 0x40, 0x40),
                   font_size=11, font_color=LIGHT50)


# ════════════════════════════════════════════
# SLIDE 8: IMPACTO
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, "IMPACTO Y ESCALABILIDAD", Inches(0), Inches(0.3), W, Inches(0.4),
         font_size=12, color=LIGHT50, alignment=PP_ALIGN.CENTER)

# Big stats
impact_stats = [
    ("2B", "usuarios WhatsApp\nen el mundo", HOPE),
    ("0€", "coste para\nel usuario", TRUST),
    ("<30s", "respuesta\ncon fuentes", WARMTH),
    ("24/7", "disponible\nsiempre", ALERT),
]
for i, (num, label, color) in enumerate(impact_stats):
    x = Inches(0.5 + i * 3.2)
    # Card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, Inches(1.2), Inches(2.8), Inches(2))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x15, 0x15, 0x15)
    card.line.color.rgb = RGBColor(0x25, 0x25, 0x25)
    card.line.width = Pt(1)

    add_text(slide, num, x + Inches(0.2), Inches(1.4), Inches(2.4), Inches(0.8),
             font_size=48, color=color, alignment=PP_ALIGN.CENTER, font_name="Georgia")
    add_text(slide, label, x + Inches(0.2), Inches(2.3), Inches(2.4), Inches(0.6),
             font_size=13, color=LIGHT50, alignment=PP_ALIGN.CENTER)

# Roadmap
roadmap = [
    ("Hoy", "IMV, empadronamiento,\nsanitaria, NIE/TIE", HOPE),
    ("Manana", "Todas las CCAA\nde Espana", TRUST),
    ("Futuro", "LATAM, refugiados,\ncualquier idioma", WARMTH),
]
for i, (phase, desc, color) in enumerate(roadmap):
    x = Inches(1.5 + i * 3.5)
    add_text(slide, phase, x, Inches(3.8), Inches(2.5), Inches(0.5),
             font_size=20, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, desc, x, Inches(4.3), Inches(2.5), Inches(0.6),
             font_size=13, color=LIGHT50, alignment=PP_ALIGN.CENTER)

    if i < 2:
        add_text(slide, "→", Inches(3.5 + i * 3.5), Inches(4), Inches(0.5), Inches(0.5),
                 font_size=24, color=LIGHT50, alignment=PP_ALIGN.CENTER)

# Purple cow / category design touch
add_text(slide, "Clara no compite con chatbots. Clara crea una nueva categoria: Civic Tenderness.",
         Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
         font_size=16, color=WARMTH, italic=True, alignment=PP_ALIGN.CENTER)

add_text(slide, "Clara no reemplaza a los trabajadores sociales. Los amplifica.",
         Inches(1.5), Inches(6.1), Inches(10), Inches(0.5),
         font_size=14, color=LIGHT50, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════
# SLIDE 9: USOS PARA VULNERABLES
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, "USOS PARA DIFERENTES PERSONAS VULNERABLES",
         Inches(0), Inches(0.3), W, Inches(0.4),
         font_size=12, color=LIGHT50, alignment=PP_ALIGN.CENTER)

use_cases = [
    ("Mayores solos", "IMV, pensiones,\nteleasistencia", WARMTH, "Voz"),
    ("Inmigrantes", "Empadronamiento,\nNIE, reagrupacion", TRUST, "Bilingue"),
    ("Violencia\nde genero", "016, ordenes\nproteccion, ayudas", ALERT, "Seguro"),
    ("Sin hogar", "Comedores, albergues,\nrenta social", HOPE, "WhatsApp"),
    ("Discapacidad", "Certificado, grado,\nayudas tecnicas", TRUST, "Audio"),
    ("Analfabetismo", "Cualquier tramite,\npaso a paso oral", WARMTH, "Voz 100%"),
]

for i, (name, desc, color, mode) in enumerate(use_cases):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.2)
    y = Inches(1.2 + row * 3)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, y, Inches(3.6), Inches(2.5))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x15, 0x15, 0x15)
    card.line.color.rgb = color
    card.line.width = Pt(1.5)

    add_text(slide, name, x + Inches(0.3), y + Inches(0.3), Inches(3), Inches(0.6),
             font_size=18, color=WHITE, bold=True)
    add_text(slide, desc, x + Inches(0.3), y + Inches(1), Inches(3), Inches(0.7),
             font_size=13, color=LIGHT50)

    # Mode tag
    add_shape_text(slide, mode, x + Inches(0.3), y + Inches(1.9), Inches(1.2), Inches(0.35),
                   fill_color=RGBColor(0x20, 0x20, 0x20), border_color=color,
                   font_size=10, font_color=color, bold=True)


# ════════════════════════════════════════════
# SLIDE 10: CIERRE
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# Team
team = [
    ("R", "Robert", "Backend", TRUST),
    ("M", "Marcos", "Deploy", HOPE),
    ("L", "Lucas", "KB + Tests", WARMTH),
    ("D", "Daniel", "Video", TRUST),
    ("A", "Andrea", "Coord", WARMTH),
]
for i, (initial, name, role, color) in enumerate(team):
    x = Inches(3 + i * 1.5)
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  x, Inches(1.5), Inches(0.7), Inches(0.7))
    circ.fill.solid()
    circ.fill.fore_color.rgb = RGBColor(0x20, 0x20, 0x20)
    circ.line.color.rgb = color
    circ.line.width = Pt(1.5)

    tf = circ.text_frame
    tf.paragraphs[0].text = initial
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = color
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text(slide, name, x - Inches(0.15), Inches(2.3), Inches(1), Inches(0.3),
             font_size=11, color=LIGHT50, alignment=PP_ALIGN.CENTER)
    add_text(slide, role, x - Inches(0.15), Inches(2.6), Inches(1), Inches(0.3),
             font_size=9, color=RGBColor(0x60, 0x60, 0x60), alignment=PP_ALIGN.CENTER)

# Big quote
add_text(slide, "Tu voz tiene poder.", Inches(0), Inches(3.5), W, Inches(1.5),
         font_size=72, color=WHITE, italic=True, alignment=PP_ALIGN.CENTER,
         font_name="Georgia")

# Warmth on "poder"
add_text(slide, "Clara", Inches(0), Inches(5.2), W, Inches(0.5),
         font_size=18, color=TRUST, alignment=PP_ALIGN.CENTER, font_name="Georgia")

add_text(slide, "OdiseIA4Good  ·  UDIT  ·  Febrero 2026",
         Inches(0), Inches(6.5), W, Inches(0.4),
         font_size=12, color=LIGHT50, alignment=PP_ALIGN.CENTER)


# ── SAVE ──
out_path = os.path.join(os.path.dirname(__file__), "Clara-Presentacion.pptx")
prs.save(out_path)
print(f"PPTX saved to: {out_path}")
print(f"Slides: {len(prs.slides)}")
